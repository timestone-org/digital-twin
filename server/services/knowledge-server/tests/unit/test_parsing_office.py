"""工作簿与演示文稿的解析。用真库现造文档，不放二进制夹具进仓。

Word 那一路在 `test_parsing_word.py`。

⚠ 现造而不是放夹具文件：一份 .docx 是一个 zip，改一个字都看不出 diff，
而「夹具是怎么来的」半年后没人说得清。现造的代价是每条用例慢几毫秒。
"""

from io import BytesIO

from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

from knowledge_server.apps.knowledge.services.parsing import RawItem
from knowledge_server.apps.knowledge.services.parsing.office import (
    PptxParser,
    XlsxParser,
)


def _raw(name: str, body: bytes) -> RawItem:
    return RawItem(filename=name, media_type="", content=body)


def _xlsx(sheets: dict[str, list[list[object]]]) -> bytes:
    book = Workbook()
    book.remove(book.active)
    for name, rows in sheets.items():
        sheet = book.create_sheet(name)
        for row in rows:
            sheet.append(row)
    buffer = BytesIO()
    book.save(buffer)
    return buffer.getvalue()


def test_xlsx_header_is_folded_into_every_row() -> None:
    """⚠ 只存 `12.5 | 开 | 3` 的话，检索到这一行也读不出它是什么——
    列名在第一行，而那一行是另一个块。"""
    made = XlsxParser().parse(
        _raw(
            "台账.xlsx",
            _xlsx({"1月": [["点位", "值"], ["出口温度", 65]]}),
        )
    )
    rows = [one for one in made.blocks if one.kind == "table_row"]
    assert rows[0].text == "点位=出口温度 | 值=65"
    assert rows[0].locator.sheet == "1月"
    assert rows[0].locator.row == 2


def test_xlsx_columns_without_a_header_still_get_kept() -> None:
    """⚠ 丢掉的话那几列的数据就再也检索不到了。"""
    made = XlsxParser().parse(
        _raw("a.xlsx", _xlsx({"表": [["点位"], ["出口温度", 65]]}))
    )
    rows = [one for one in made.blocks if one.kind == "table_row"]
    assert "第2列=65" in rows[0].text


def test_xlsx_blank_rows_are_dropped() -> None:
    made = XlsxParser().parse(
        _raw(
            "a.xlsx",
            _xlsx({"表": [["点位"], [None, None], ["出口温度"]]}),
        )
    )
    rows = [one for one in made.blocks if one.kind == "table_row"]
    assert len(rows) == 1


def test_xlsx_scans_every_sheet() -> None:
    made = XlsxParser().parse(
        _raw(
            "a.xlsx",
            _xlsx({"甲": [["列"], ["一"]], "乙": [["列"], ["二"]]}),
        )
    )
    sheets = {one.locator.sheet for one in made.blocks}
    assert sheets == {"甲", "乙"}


def _pptx(slides: list[list[str]]) -> bytes:
    deck = Presentation()
    blank = deck.slide_layouts[6]
    for texts in slides:
        slide = deck.slides.add_slide(blank)
        for index, text in enumerate(texts):
            box = slide.shapes.add_textbox(
                Inches(1), Inches(1 + index), Inches(4), Inches(1)
            )
            box.text_frame.text = text
    buffer = BytesIO()
    deck.save(buffer)
    return buffer.getvalue()


def test_pptx_first_box_on_a_slide_is_its_heading() -> None:
    made = PptxParser().parse(
        _raw("汇报.pptx", _pptx([["冷却水系统", "出口温度 65 ℃"]]))
    )
    assert made.blocks[0].kind == "heading"
    assert made.blocks[1].locator.path == ("冷却水系统",)


def test_pptx_page_numbers_start_at_one() -> None:
    made = PptxParser().parse(_raw("a.pptx", _pptx([["甲"], ["乙"]])))
    assert [one.locator.page for one in made.blocks] == [1, 2]


def test_pptx_shapes_without_text_never_raise() -> None:
    """⚠ 取一个图片形状的 `text` 会抛，而一次抛就是整份文档摄取失败。"""
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(4), Inches(1))
    box.text_frame.text = "正文"
    buffer = BytesIO()
    deck.save(buffer)
    made = PptxParser().parse(_raw("a.pptx", buffer.getvalue()))
    assert [one.text for one in made.blocks] == ["正文"]
