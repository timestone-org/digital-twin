"""Office 三件套：docx / xlsx / pptx。

⚠ 三种格式合在一个文件里，是因为它们**共用同一条结构提取思路**（标题层级 /
工作表与行 / 幻灯片序号），拆成三个文件会让那条思路被抄三遍。加第四种 Office
格式时再拆。

⚠ 这三路的解析都是**纯 CPU 且阻塞**的。调用方必须扔进进程池——放进事件循环
会把整条消费循环连同健康探针一起冻住，而现象是「服务好好的，队列不动了」。

⚠ 都只读**文本与结构**：不取嵌入的图片、不跑宏、不跟外部引用。一份 Office
文档里可以有任何东西。
"""

from dataclasses import dataclass
from io import BytesIO
from typing import cast

from docx import Document
from docx.document import Document as DocxDocument
from openpyxl import load_workbook
from openpyxl.worksheet.worksheet import Worksheet
from pptx import Presentation
from pptx.shapes.autoshape import Shape
from pptx.slide import Slide

from knowledge_server.apps.knowledge.services.parsing.ports import (
    Block,
    Locator,
    ParsedDocument,
    RawItem,
)
from knowledge_server.apps.knowledge.services.parsing.structure import (
    cell_text,
    paired,
    pushed,
)

# 与纯文本那一路同一个上限，理由也同源
MAX_BLOCKS = 20_000

_DOCX_SUFFIXES = (".docx",)
_XLSX_SUFFIXES = (".xlsx", ".xlsm")
_PPTX_SUFFIXES = (".pptx",)


def _heading_level(style_name: str) -> int:
    """docx 的样式名 → 标题层级；不是标题给 0。

    ⚠ 认 `Heading N` 也认中文的「标题 N」：现场的模板多半是中文版 Word 存的，
    只认英文的话整份文档一个标题都解不出来，而那正是切块质量的主要来源。

    Args: style_name。
    """
    for prefix in ("Heading ", "标题 ", "标题"):
        if style_name.startswith(prefix):
            tail = style_name[len(prefix) :].strip()
            if tail.isdigit():
                return int(tail)
    return 0


@dataclass(frozen=True)
class DocxParser:
    """Word 文档：按段落走，`Heading N` 撑起标题层级。"""

    name: str = "docx"
    suffixes: tuple[str, ...] = _DOCX_SUFFIXES
    media_types: tuple[str, ...] = (
        "application/vnd.openxmlformats-officedocument"
        ".wordprocessingml.document",
    )

    def parse(self, raw: RawItem) -> ParsedDocument:
        """段落与表格行都收，标题路径一路带下去。

        Args: raw。
        """
        document = Document(BytesIO(raw.content))
        stack: list[tuple[int, str]] = []
        made: list[Block] = []
        for paragraph in document.paragraphs:
            if len(made) >= MAX_BLOCKS:
                break
            text = paragraph.text.strip()
            if not text:
                continue
            style = paragraph.style
            level = _heading_level("" if style is None else (style.name or ""))
            if level > 0:
                stack = pushed(stack, level, text)
            path = tuple(one for _level, one in stack)
            made.append(
                Block(
                    kind="heading" if level > 0 else "paragraph",
                    text=text,
                    level=level,
                    locator=Locator(path=path),
                )
            )
        made.extend(_docx_tables(document, stack, len(made)))
        return ParsedDocument(
            title=raw.filename,
            blocks=tuple(made),
            is_truncated=len(made) >= MAX_BLOCKS,
        )


def _docx_tables(
    document: DocxDocument, stack: list[tuple[int, str]], made_so_far: int
) -> list[Block]:
    """把文档里的表格摊成竖线行。

    ⚠ 表格排在段落之后而不是原位：python-docx 的 `paragraphs` 与 `tables`
    是两份平行清单，拿不到它们在正文里的真实先后。摆在后面至少顺序是确定的，
    而穿插着猜会让引用指错地方。

    Args: document, stack, made_so_far。
    """
    path = tuple(one for _level, one in stack)
    made: list[Block] = []
    tables = document.tables
    for index, table in enumerate(tables, start=1):
        for row_index, row in enumerate(table.rows, start=1):
            if made_so_far + len(made) >= MAX_BLOCKS:
                return made
            cells = [cell_text(one.text) for one in row.cells]
            if not any(cells):
                continue
            made.append(
                Block(
                    kind="table_row",
                    text=" | ".join(cells),
                    locator=Locator(
                        sheet=f"表 {index}", row=row_index, path=path
                    ),
                )
            )
    return made


@dataclass(frozen=True)
class XlsxParser:
    """工作簿：一行一块，表头拼进正文。

    ⚠ **表头要拼进每一行**：只存 `12.5 | 开 | 3` 的话，检索到这一行也读不出
    它是什么——列名在第一行，而那一行是另一个块。拼进去之后每一块都自足。

    ⚠ 只读值不读公式（`data_only=True`）：读公式串回来的是 `=SUM(A1:A9)`，
    检索它毫无意义。代价是没打开过的表取到 `None`——那时如实当空。
    """

    name: str = "xlsx"
    suffixes: tuple[str, ...] = _XLSX_SUFFIXES
    media_types: tuple[str, ...] = (
        "application/vnd.openxmlformats-officedocument" ".spreadsheetml.sheet",
    )

    def parse(self, raw: RawItem) -> ParsedDocument:
        """逐个工作表逐行摊平。

        Args: raw。
        """
        book = load_workbook(
            BytesIO(raw.content), data_only=True, read_only=True
        )
        made: list[Block] = []
        try:
            for sheet in book.worksheets:
                made.extend(_sheet_blocks(sheet, len(made)))
                if len(made) >= MAX_BLOCKS:
                    break
        finally:
            # ⚠ read_only 模式握着文件句柄，不关会攒着不放
            book.close()
        return ParsedDocument(
            title=raw.filename,
            blocks=tuple(made),
            is_truncated=len(made) >= MAX_BLOCKS,
        )


def _sheet_blocks(sheet: Worksheet, made_so_far: int) -> list[Block]:
    """一个工作表摊成行块，第一行当表头。

    Args: sheet, made_so_far。
    """
    title = str(sheet.title)
    header: list[str] = []
    made: list[Block] = []
    for row_index, row in enumerate(sheet.iter_rows(values_only=True), 1):
        if made_so_far + len(made) >= MAX_BLOCKS:
            break
        cells = [cell_text(one) for one in row]
        if not any(cells):
            continue
        if not header:
            header = cells
            made.append(
                Block(
                    kind="heading",
                    text=" | ".join(cells),
                    level=1,
                    locator=Locator(sheet=title, row=row_index),
                )
            )
            continue
        made.append(
            Block(
                kind="table_row",
                text=paired(header, cells),
                locator=Locator(sheet=title, row=row_index, path=(title,)),
            )
        )
    return made


@dataclass(frozen=True)
class PptxParser:
    """演示文稿：一页一组块，页码进 locator。"""

    name: str = "pptx"
    suffixes: tuple[str, ...] = _PPTX_SUFFIXES
    media_types: tuple[str, ...] = (
        "application/vnd.openxmlformats-officedocument"
        ".presentationml.presentation",
    )

    def parse(self, raw: RawItem) -> ParsedDocument:
        """逐页取文本框，每页第一个当这页的标题。

        Args: raw。
        """
        deck = Presentation(BytesIO(raw.content))
        made: list[Block] = []
        for page, slide in enumerate(deck.slides, start=1):
            made.extend(_slide_blocks(slide, page, len(made)))
            if len(made) >= MAX_BLOCKS:
                break
        return ParsedDocument(
            title=raw.filename,
            blocks=tuple(made),
            is_truncated=len(made) >= MAX_BLOCKS,
        )


def _slide_blocks(slide: Slide, page: int, made_so_far: int) -> list[Block]:
    """一页里的文本框。第一个非空的当这页标题。

    Args: slide, page, made_so_far。
    """
    made: list[Block] = []
    heading = ""
    for shape in slide.shapes:
        if made_so_far + len(made) >= MAX_BLOCKS:
            break
        # ⚠ 没有 `text_frame` 的形状（图片、连线、表格容器）直接跳过：
        # 取它的 `text` 会抛，而一次抛就是整份文档摄取失败
        if not shape.has_text_frame:
            continue
        # ⚠ `has_text_frame` 挂在 `BaseShape` 上，`text_frame` 只在带文本的
        # 子类上——上一行判过之后收窄一次，不然类型检查看到的是「基类没有
        # 这个属性」，而运行期它是好的
        text = cast("Shape", shape).text_frame.text.strip()
        if not text:
            continue
        if not heading:
            heading = text
            made.append(
                Block(
                    kind="heading",
                    text=text,
                    level=1,
                    locator=Locator(page=page, path=(text,)),
                )
            )
            continue
        made.append(
            Block(
                kind="paragraph",
                text=text,
                locator=Locator(page=page, path=(heading,)),
            )
        )
    return made
